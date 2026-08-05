#!/usr/bin/env python3
"""Create a standalone slide explaining M2 missing-value imputation.

Outputs:
    presentation/m2_missing_values_imputation_slide.pptx
    presentation/m2_missing_values_imputation_slide.pdf
"""
from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
GREY = RGBColor(0x55, 0x5B, 0x66)
ACCENT = RGBColor(0x1F, 0x6F, 0xB2)
LIGHT = RGBColor(0xEC, 0xF1, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF5, 0xF8, 0xFB)
AMBER = RGBColor(0xB7, 0x80, 0x1E)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
NEG = RGBColor(0xC0, 0x3A, 0x2B)

OUT_STEM = "m2_missing_values_imputation_slide"
FOOT = "SSGA Field Project - Final Presentation"


def _tb(slide, left, top, width, height, anchor=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def _rect(slide, left, top, width, height, color, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _para(tf, text, size, color, bold=False, first=False, before=0, after=4, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    if align:
        p.alignment = align
    return p


def _header(slide):
    _rect(slide, 0.0, 0.0, 0.22, 7.5, NAVY)
    k = _tb(slide, 0.62, 0.30, 12.4, 0.4)
    _para(k, "PART 2 - M2 PREPROCESSING", 12, ACCENT, bold=True, first=True, after=0)
    t = _tb(slide, 0.6, 0.66, 12.4, 0.70)
    _para(t, "Which Missing Values Are Imputed Before M2?", 25, NAVY, bold=True, first=True, after=0)
    _rect(slide, 0.62, 1.42, 4.2, 0.045, ACCENT)


def _footer(slide):
    f = _tb(slide, 0.6, 7.05, 12.0, 0.35)
    _para(f, FOOT, 9, GREY, first=True, after=0)


def _pill(slide, left, top, width, label, fill, txt_color=WHITE):
    _rect(slide, left, top, width, 0.36, fill, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = _tb(slide, left + 0.04, top + 0.06, width - 0.08, 0.22)
    _para(tf, label, 9, txt_color, bold=True, first=True, after=0, align=PP_ALIGN.CENTER)


def _card(slide, left, top, width, height, title, lines, title_color=ACCENT):
    _rect(slide, left, top, width, height, CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = _tb(slide, left + 0.20, top + 0.14, width - 0.40, height - 0.25)
    _para(tf, title.upper(), 10.8, title_color, bold=True, first=True, after=5)
    for i, line in enumerate(lines):
        _para(tf, line, 10.9 if i else 11.3, NAVY if i == 0 else GREY, bold=(i == 0), after=2)


def build_pptx(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide)

    intro = _tb(slide, 0.72, 1.62, 12.0, 0.42)
    _para(
        intro,
        "In code: M2 runs SimpleImputer(strategy='median') before StandardScaler, logistic regression, and sigmoid calibration.",
        12.3,
        GREY,
        first=True,
        after=0,
    )

    # Pipeline strip.
    x0, y0 = 0.78, 2.15
    steps = [
        ("1  Fit rows", "M1_signal != 0 and meta_label exists", ACCENT),
        ("2  Impute", "each feature's training median", AMBER),
        ("3  Scale", "mean 0, std 1", GREEN),
        ("4  Probability", "calibrated p_success", NAVY),
    ]
    for i, (label, sub, color) in enumerate(steps):
        left = x0 + i * 3.05
        _rect(slide, left, y0, 2.62, 0.78, color, MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = _tb(slide, left + 0.12, y0 + 0.10, 2.38, 0.58)
        _para(tf, label, 12.0, WHITE, bold=True, first=True, after=1)
        _para(tf, sub, 8.8, WHITE, after=0)
        if i < len(steps) - 1:
            ar = _tb(slide, left + 2.66, y0 + 0.16, 0.34, 0.36, MSO_ANCHOR.MIDDLE)
            _para(ar, "->", 15, ACCENT, bold=True, first=True, after=0, align=PP_ALIGN.CENTER)

    _card(
        slide,
        0.70,
        3.18,
        3.88,
        2.27,
        "Actual NaNs in train fit rows",
        [
            "23 of 52 M2 columns had missing values.",
            "Mostly rolling lookbacks: mom_4w/12w/26w/52w, vol_4w/12w/26w, trend_signal, drawdown_26w, corr_to_spy_26w.",
            "Also derived relatives: z_* versions, rel_mom_12w, mom_vol_interaction, dispersion_4w/12w, avg_pairwise_corr_26w, m1_x_vol.",
            "Counts are largest for 52w features because they need the longest history.",
        ],
        NEG,
    )

    _card(
        slide,
        4.82,
        3.18,
        3.64,
        2.27,
        "How the fill works",
        [
            "Median is learned feature-by-feature from training rows only.",
            "Example: missing mom_52w -> training median of mom_52w, not zero and not a future value.",
            "The same fitted medians are reused at prediction time.",
            "Labels, future returns, p_success, and raw prices are not M2 inputs, so they are not imputed here.",
        ],
        ACCENT,
    )

    _card(
        slide,
        8.70,
        3.18,
        3.94,
        2.27,
        "Why it happens",
        [
            "Rolling features have a warmup period.",
            "A 52w momentum needs 52 prior weekly closes, then shift(1) adds one more no-lookahead delay.",
            "Cross-sectional z-scores can be missing if all sleeves look the same that week, e.g. drawdown = 0 for each selected sleeve.",
            "Macro/VIX and asset-class dummy columns were complete in the saved training rows.",
        ],
        GREEN,
    )

    _rect(slide, 0.72, 5.70, 11.92, 0.34, WHITE)
    note = _tb(slide, 0.82, 5.74, 11.70, 0.22)
    _para(
        note,
        "Saved long-only run: 1,500 M2 train rows; top missing counts were mom_52w / z_mom_52w = 159, trend / z_trend = 120, and 26w risk/correlation features = 81.",
        10.3,
        GREY,
        first=True,
        after=0,
    )

    _rect(slide, 0.60, 6.30, 12.13, 0.56, LIGHT)
    tk = _tb(slide, 0.84, 6.34, 11.70, 0.48, MSO_ANCHOR.MIDDLE)
    _para(
        tk,
        "Takeaway: imputation is a defensive preprocessing step for early-window and zero-dispersion gaps; it does not manufacture labels or use future information.",
        12.2,
        ACCENT,
        bold=True,
        first=True,
        after=0,
    )
    _footer(slide)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def build_pdf(out_path: Path) -> None:
    fig = plt.figure(figsize=(13.333, 7.5), dpi=72)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 13.333)
    ax.set_ylim(0, 7.5)
    ax.axis("off")

    navy = "#0B2A4A"
    grey = "#555B66"
    accent = "#1F6FB2"
    light = "#ECF1F6"
    card = "#F5F8FB"
    amber = "#B7801E"
    green = "#1E8E5A"
    neg = "#C03A2B"
    white = "#FFFFFF"

    ax.add_patch(Rectangle((0, 0), 0.22, 7.5, color=navy))
    ax.text(0.62, 7.10, "PART 2 - M2 PREPROCESSING", fontsize=12, color=accent, fontweight="bold", va="top")
    ax.text(0.60, 6.84, "Which Missing Values Are Imputed Before M2?", fontsize=25, color=navy, fontweight="bold", va="top")
    ax.add_patch(Rectangle((0.62, 6.03), 4.2, 0.045, color=accent))
    ax.text(
        0.72,
        5.76,
        "In code: M2 runs SimpleImputer(strategy='median') before StandardScaler, logistic regression, and sigmoid calibration.",
        fontsize=12.3,
        color=grey,
        va="top",
    )

    steps = [
        ("1  Fit rows", "M1_signal != 0 and meta_label exists", accent),
        ("2  Impute", "each feature's training median", amber),
        ("3  Scale", "mean 0, std 1", green),
        ("4  Probability", "calibrated p_success", navy),
    ]
    for i, (label, sub, color) in enumerate(steps):
        x = 0.78 + i * 3.05
        ax.add_patch(Rectangle((x, 4.57), 2.62, 0.78, color=color))
        ax.text(x + 0.12, 5.17, label, fontsize=12, color=white, fontweight="bold", va="top")
        ax.text(x + 0.12, 4.89, sub, fontsize=8.8, color=white, va="top")
        if i < 3:
            ax.text(x + 2.72, 4.96, "->", fontsize=15, color=accent, fontweight="bold", va="center")

    cards = [
        (
            0.70,
            "ACTUAL NaNs IN TRAIN FIT ROWS",
            neg,
            [
                "23 of 52 M2 columns had missing values.",
                "Mostly rolling lookbacks: mom_4w/12w/26w/52w, vol_4w/12w/26w, trend_signal, drawdown_26w, corr_to_spy_26w.",
                "Also derived relatives: z_* versions, rel_mom_12w, mom_vol_interaction, dispersion_4w/12w, avg_pairwise_corr_26w, m1_x_vol.",
                "Counts are largest for 52w features because they need the longest history.",
            ],
        ),
        (
            4.82,
            "HOW THE FILL WORKS",
            accent,
            [
                "Median is learned feature-by-feature from training rows only.",
                "Example: missing mom_52w -> training median of mom_52w, not zero and not a future value.",
                "The same fitted medians are reused at prediction time.",
                "Labels, future returns, p_success, and raw prices are not M2 inputs, so they are not imputed here.",
            ],
        ),
        (
            8.70,
            "WHY IT HAPPENS",
            green,
            [
                "Rolling features have a warmup period.",
                "A 52w momentum needs 52 prior weekly closes, then shift(1) adds one more no-lookahead delay.",
                "Cross-sectional z-scores can be missing if all sleeves look the same that week, e.g. drawdown = 0 for each selected sleeve.",
                "Macro/VIX and asset-class dummy columns were complete in the saved training rows.",
            ],
        ),
    ]
    widths = [3.88, 3.64, 3.94]
    for (x, title, color, lines), width in zip(cards, widths):
        ax.add_patch(Rectangle((x, 2.05), width, 2.27, color=card))
        ax.text(x + 0.20, 4.15, title, fontsize=10.4, color=color, fontweight="bold", va="top")
        y = 3.78
        for j, line in enumerate(lines):
            ax.text(x + 0.20, y, line, fontsize=9.35 if len(line) > 95 else 10.2, color=navy if j == 0 else grey, fontweight="bold" if j == 0 else "normal", va="top", wrap=True)
            y -= 0.40

    ax.text(
        0.82,
        1.70,
        "Saved long-only run: 1,500 M2 train rows; top missing counts were mom_52w / z_mom_52w = 159, trend / z_trend = 120, and 26w risk/correlation features = 81.",
        fontsize=10.3,
        color=grey,
        va="top",
    )
    ax.add_patch(Rectangle((0.60, 0.64), 12.13, 0.56, color=light))
    ax.text(
        0.84,
        1.03,
        "Takeaway: imputation is a defensive preprocessing step for early-window and zero-dispersion gaps; it does not manufacture labels or use future information.",
        fontsize=12.2,
        color=accent,
        fontweight="bold",
        va="top",
    )
    ax.text(0.60, 0.30, FOOT, fontsize=9, color=grey, va="bottom")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, format="pdf", facecolor="white")
    plt.close(fig)


def main() -> None:
    root = Path(__file__).resolve().parent.parent
    pptx_out = root / "presentation" / f"{OUT_STEM}.pptx"
    pdf_out = root / "presentation" / f"{OUT_STEM}.pdf"
    build_pptx(pptx_out)
    build_pdf(pdf_out)
    print(f"saved {pptx_out}")
    print(f"saved {pdf_out}")


if __name__ == "__main__":
    main()
