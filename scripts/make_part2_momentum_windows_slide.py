#!/usr/bin/env python3
"""Create a standalone Part 2 slide on 12w/26w/52w momentum windows.

Outputs:
    presentation/part2_momentum_windows_variability_slide.pptx
    presentation/part2_momentum_windows_variability_slide.pdf

The slide matches the visual language used in the final deck:
16:9 canvas, navy rail, blue accent, compact explanatory cards, and a
takeaway strip.
"""
from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
GREY = RGBColor(0x55, 0x5B, 0x66)
ACCENT = RGBColor(0x1F, 0x6F, 0xB2)
LIGHT = RGBColor(0xEC, 0xF1, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF5, 0xF8, 0xFB)
AMBER = RGBColor(0xB7, 0x80, 0x1E)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)

OUT_STEM = "part2_momentum_windows_variability_slide"
FOOT = "SSGA Field Project - Final Presentation"


def _tb(slide, left, top, width, height, anchor=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def _rect(slide, left, top, width, height, color, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _para(tf, text, size, color, bold=False, first=False, before=0, after=4, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    if align:
        p.alignment = align
    return p


def _header(slide):
    _rect(slide, 0.0, 0.0, 0.22, 7.5, NAVY)
    k = _tb(slide, 0.62, 0.30, 12.4, 0.4)
    _para(k, "PART 2 - M1 MOMENTUM", 12, ACCENT, bold=True, first=True, after=0)
    t = _tb(slide, 0.6, 0.66, 12.4, 0.70)
    _para(t, "Momentum Windows: Variability and Why 12 / 26 / 52 Weeks", 24, NAVY, bold=True, first=True, after=0)
    _rect(slide, 0.62, 1.42, 4.2, 0.045, ACCENT)


def _footer(slide):
    f = _tb(slide, 0.6, 7.05, 12.0, 0.35)
    _para(f, FOOT, 9, GREY, first=True, after=0)


def _card(slide, left, top, width, height, label, headline, bullets, label_color=ACCENT):
    _rect(slide, left, top, width, height, CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = _tb(slide, left + 0.20, top + 0.15, width - 0.40, height - 0.30)
    _para(tf, label.upper(), 10.5, label_color, bold=True, first=True, after=3)
    _para(tf, headline, 15, NAVY, bold=True, after=5)
    for bullet in bullets:
        _para(tf, "- " + bullet, 10.9, GREY, after=2)


def _bar(slide, left, top, label, values, colors):
    title = _tb(slide, left, top, 2.95, 0.28)
    _para(title, label, 10.5, NAVY, bold=True, first=True, after=0)
    names = ["12w", "26w", "52w"]
    max_val = max(values)
    for i, (name, val, col) in enumerate(zip(names, values, colors)):
        y = top + 0.38 + i * 0.28
        nm = _tb(slide, left, y - 0.03, 0.36, 0.18)
        _para(nm, name, 7.7, GREY, first=True, after=0)
        _rect(slide, left + 0.42, y, 1.65 * val / max_val, 0.11, col)


def build_pptx(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide)

    eq = _tb(slide, 0.72, 1.60, 7.05, 0.45)
    _para(
        eq,
        "Formula: momentum_N = price_today / price_N_weeks_ago - 1, then z-scored cross-sectionally across sleeves.",
        12.2,
        GREY,
        first=True,
        after=0,
    )

    # Window cards.
    _card(
        slide,
        0.70,
        2.10,
        3.78,
        2.45,
        "12 weeks",
        "Recent pulse",
        [
            "Most reactive week to week; lowest raw accumulated spread.",
            "Catches earnings surprises, CPI/Fed repricing, sudden risk-on/off moves.",
            "Useful when the market has just changed its mind.",
        ],
        ACCENT,
    )
    _card(
        slide,
        4.78,
        2.10,
        3.78,
        2.45,
        "26 weeks",
        "Regime bridge",
        [
            "Middle variability; filters one noisy quarter without becoming stale.",
            "Spans half a year: rate-cycle shifts, credit tightening, summer/winter demand.",
            "Balances timeliness with persistence.",
        ],
        AMBER,
    )
    _card(
        slide,
        8.86,
        2.10,
        3.78,
        2.45,
        "52 weeks",
        "Full-year trend",
        [
            "Highest raw dispersion because a year compounds more shocks.",
            "Covers four earnings seasons, holiday retail, winter energy, tax/year-end flows.",
            "Helps avoid mistaking a seasonal burst for a durable winner.",
        ],
        GREEN,
    )

    # Comparison panel.
    _rect(slide, 0.70, 4.90, 5.95, 1.14, LIGHT)
    c1 = _tb(slide, 0.92, 4.99, 5.55, 0.22)
    _para(c1, "Which is more variable?", 12.0, NAVY, bold=True, first=True, after=0)
    c2 = _tb(slide, 0.92, 5.27, 5.55, 0.58)
    _para(
        c2,
        "For raw N-week momentum, 52w > 26w > 12w in total variation: longer windows accumulate more return shocks and let winners/losers separate further. But 12w is the most reactive.",
        11.2,
        GREY,
        first=True,
        after=0,
    )

    _rect(slide, 6.95, 4.90, 5.69, 1.14, WHITE)
    _bar(slide, 7.18, 5.03, "Raw momentum dispersion", [1, 1.45, 2.05], [ACCENT, AMBER, GREEN])
    _bar(slide, 9.85, 5.03, "Week-to-week reactivity", [2.05, 1.45, 1], [ACCENT, AMBER, GREEN])

    _rect(slide, 0.60, 6.28, 12.13, 0.58, LIGHT)
    tk = _tb(slide, 0.84, 6.33, 11.70, 0.50, MSO_ANCHOR.MIDDLE)
    _para(
        tk,
        "Takeaway: use all three because they answer different questions - what changed recently, what persisted for half a year, and what survived a full seasonal cycle.",
        12.4,
        ACCENT,
        bold=True,
        first=True,
        after=0,
    )
    _footer(slide)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def build_pdf(out_path: Path) -> None:
    # Matplotlib duplicate for easy viewing when PowerPoint export is unavailable.
    fig = plt.figure(figsize=(13.333, 7.5), dpi=72)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 13.333)
    ax.set_ylim(0, 7.5)
    ax.axis("off")

    navy = "#0B2A4A"
    grey = "#555B66"
    accent = "#1F6FB2"
    light = "#ECF1F6"
    card = "#F5F8FB"
    amber = "#B7801E"
    green = "#1E8E5A"

    ax.add_patch(Rectangle((0, 0), 0.22, 7.5, color=navy))
    ax.text(0.62, 7.10, "PART 2 - M1 MOMENTUM", fontsize=12, color=accent, fontweight="bold", va="top")
    ax.text(
        0.60,
        6.84,
        "Momentum Windows: Variability and Why 12 / 26 / 52 Weeks",
        fontsize=24,
        color=navy,
        fontweight="bold",
        va="top",
    )
    ax.add_patch(Rectangle((0.62, 6.03), 4.2, 0.045, color=accent))
    ax.text(
        0.72,
        5.74,
        "Formula: momentum_N = price_today / price_N_weeks_ago - 1, then z-scored cross-sectionally across sleeves.",
        fontsize=12.2,
        color=grey,
        va="top",
    )

    cards = [
        (
            0.70,
            "12 WEEKS",
            "Recent pulse",
            accent,
            [
                "Most reactive week to week; lowest raw accumulated spread.",
                "Catches earnings surprises, CPI/Fed repricing, sudden risk-on/off moves.",
                "Useful when the market has just changed its mind.",
            ],
        ),
        (
            4.78,
            "26 WEEKS",
            "Regime bridge",
            amber,
            [
                "Middle variability; filters one noisy quarter without becoming stale.",
                "Spans half a year: rate-cycle shifts, credit tightening, summer/winter demand.",
                "Balances timeliness with persistence.",
            ],
        ),
        (
            8.86,
            "52 WEEKS",
            "Full-year trend",
            green,
            [
                "Highest raw dispersion because a year compounds more shocks.",
                "Covers four earnings seasons, holiday retail, winter energy, tax/year-end flows.",
                "Helps avoid mistaking a seasonal burst for a durable winner.",
            ],
        ),
    ]
    for x, label, headline, color, bullets in cards:
        ax.add_patch(Rectangle((x, 2.95), 3.78, 2.45, color=card))
        ax.text(x + 0.20, 5.20, label, fontsize=10.5, color=color, fontweight="bold", va="top")
        ax.text(x + 0.20, 4.86, headline, fontsize=15, color=navy, fontweight="bold", va="top")
        y = 4.48
        for bullet in bullets:
            ax.text(x + 0.20, y, "- " + bullet, fontsize=10.4, color=grey, va="top", wrap=True)
            y -= 0.42

    ax.add_patch(Rectangle((0.70, 1.46), 5.95, 1.14, color=light))
    ax.text(0.92, 2.40, "Which is more variable?", fontsize=12, color=navy, fontweight="bold", va="top")
    ax.text(
        0.92,
        2.10,
        "For raw N-week momentum, 52w > 26w > 12w in total variation: longer windows accumulate more return shocks and let winners/losers separate further. But 12w is the most reactive.",
        fontsize=10.7,
        color=grey,
        va="top",
        wrap=True,
    )

    def bar_group(x, title, vals):
        ax.text(x, 2.37, title, fontsize=10.5, color=navy, fontweight="bold", va="top")
        max_val = max(vals)
        for i, (nm, val, col) in enumerate(zip(["12w", "26w", "52w"], vals, [accent, amber, green])):
            y = 2.03 - i * 0.28
            ax.text(x, y + 0.05, nm, fontsize=7.7, color=grey, va="center")
            ax.add_patch(Rectangle((x + 0.42, y), 1.65 * val / max_val, 0.11, color=col))

    bar_group(7.18, "Raw momentum dispersion", [1, 1.45, 2.05])
    bar_group(9.85, "Week-to-week reactivity", [2.05, 1.45, 1])

    ax.add_patch(Rectangle((0.60, 0.64), 12.13, 0.58, color=light))
    ax.text(
        0.84,
        1.03,
        "Takeaway: use all three because they answer different questions - what changed recently, what persisted for half a year, and what survived a full seasonal cycle.",
        fontsize=12.4,
        color=accent,
        fontweight="bold",
        va="top",
    )
    ax.text(0.60, 0.30, FOOT, fontsize=9, color=grey, va="bottom")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, format="pdf", facecolor="white")
    plt.close(fig)


def main() -> None:
    root = Path(__file__).resolve().parent.parent
    pptx_out = root / "presentation" / f"{OUT_STEM}.pptx"
    pdf_out = root / "presentation" / f"{OUT_STEM}.pdf"
    build_pptx(pptx_out)
    build_pdf(pdf_out)
    print(f"saved {pptx_out}")
    print(f"saved {pdf_out}")


if __name__ == "__main__":
    main()
