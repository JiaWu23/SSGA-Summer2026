#!/usr/bin/env python3
"""Generate the deliverables review deck (answers Xuesong's four asks).

    python scripts/make_deliverables_deck.py
        -> presentation/SSGA_Deliverables_Deck.pptx

House style matched to scripts/make_summer_deck.py (navy / accent, 16:9, blank layout).
Written to be *explainable*: every slide reads on its own, with full-sentence bullets and a
takeaway strip. Numbers are from reproduced walk-forward runs (see runs/*_sweep/).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
GREY = RGBColor(0x55, 0x5B, 0x66)
ACCENT = RGBColor(0x1F, 0x6F, 0xB2)
LIGHT = RGBColor(0xEC, 0xF1, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
POS = RGBColor(0x1E, 0x8E, 0x5A)
NEG = RGBColor(0xC0, 0x3A, 0x2B)
AMBER = RGBColor(0xB7, 0x80, 0x1E)
CARD = RGBColor(0xF5, 0xF8, 0xFB)

DATE = "Next update · July 2026"
FOOT = "SSGA Summer 2026 · Meta-Labeling Pipeline · Deliverables Review"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_n = [0]


def _tb(slide, left, top, width, height, anchor=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def _rect(slide, left, top, width, height, color, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _para(tf, text, size, color, bold=False, first=False, before=0, after=4, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if before:
        p.space_before = Pt(before)
    p.space_after = Pt(after)
    if align:
        p.alignment = align
    return p


def _header(slide, title, kicker=None):
    _rect(slide, 0.0, 0.0, 0.22, 7.5, NAVY)
    if kicker:
        k = _tb(slide, 0.62, 0.30, 12.4, 0.4)
        _para(k, kicker.upper(), 12, ACCENT, bold=True, first=True, after=0)
        t = _tb(slide, 0.6, 0.66, 12.4, 0.9)
    else:
        t = _tb(slide, 0.6, 0.4, 12.4, 0.9)
    _para(t, title, 26, NAVY, bold=True, first=True, after=0)
    _rect(slide, 0.62, 1.42 if kicker else 1.24, 4.2, 0.045, ACCENT)


def _footer(slide):
    _n[0] += 1
    f = _tb(slide, 0.6, 7.05, 12.0, 0.35)
    _para(f, FOOT, 9, GREY, first=True, after=0)
    pg = _tb(slide, 12.5, 7.05, 0.6, 0.35)
    _para(pg, str(_n[0]), 9, GREY, first=True, after=0)


def _card(slide, left, top, width, height, label, lines, label_color=ACCENT, fill=CARD):
    _rect(slide, left, top, width, height, fill, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = _tb(slide, left + 0.22, top + 0.16, width - 0.44, height - 0.32)
    _para(tf, label.upper(), 11, label_color, bold=True, first=True, after=6)
    for txt, sz, col, bd in lines:
        _para(tf, txt, sz, col, bold=bd, after=3)
    return tf


# ---------------------------------------------------------------- slide builders
def title_slide(title, subtitle, presenter):
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0.0, 0.0, 13.333, 0.35, NAVY)
    _rect(s, 0.0, 7.15, 13.333, 0.35, ACCENT)
    tf = _tb(s, 0.9, 2.0, 11.5, 3.2)
    _para(tf, title, 38, NAVY, bold=True, first=True, after=0)
    _para(tf, subtitle, 18, GREY, before=12, after=0)
    _para(tf, DATE, 15, ACCENT, bold=True, before=14, after=0)
    _para(tf, presenter, 13, GREY, before=6, after=0)
    return s


def bullets_slide(title, kicker, bullets, takeaway=None, takeaway_color=ACCENT):
    s = prs.slides.add_slide(BLANK)
    _header(s, title, kicker)
    body = _tb(s, 0.7, 1.7, 12.0, 4.3)
    for i, (txt, lvl) in enumerate(bullets):
        _para(body, ("■  " if lvl == 0 else "–  ") + txt,
              16 if lvl == 0 else 13.5, NAVY if lvl == 0 else GREY,
              bold=(lvl == 0), first=(i == 0), after=5)
    if takeaway:
        _rect(s, 0.6, 6.2, 12.13, 0.66, LIGHT)
        tk = _tb(s, 0.85, 6.24, 11.7, 0.58, MSO_ANCHOR.MIDDLE)
        _para(tk, "Takeaway:  " + takeaway, 13, takeaway_color, bold=True, first=True, after=0)
    _footer(s)
    return s


def table_slide(title, kicker, headers, rows, note=None, hi_row=None, hi_col_map=None):
    """rows: list of list[str]. hi_row: index (1-based data) to highlight. hi_col_map:
    dict {(row_idx, col_idx): RGBColor} for per-cell font color."""
    s = prs.slides.add_slide(BLANK)
    _header(s, title, kicker)
    nr, nc = len(rows) + 1, len(headers)
    top = 1.75
    rh = min(0.42, 4.1 / nr)
    gt = s.shapes.add_table(nr, nc, Inches(0.7), Inches(top), Inches(11.93), Inches(rh * nr)).table
    for c, h in enumerate(headers):
        cell = gt.cell(0, c)
        cell.text = h
        pr = cell.text_frame.paragraphs[0]
        pr.font.size = Pt(12); pr.font.bold = True; pr.font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.text = str(val)
            pr = cell.text_frame.paragraphs[0]
            pr.font.size = Pt(11.5)
            col = GREY
            if hi_col_map and (r, c) in hi_col_map:
                col = hi_col_map[(r, c)]; pr.font.bold = True
            elif c == 0:
                col = NAVY
            pr.font.color.rgb = col
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if (hi_row == r) else WHITE
    if note:
        nb = _tb(s, 0.7, top + rh * nr + 0.18, 11.9, 1.2)
        _para(nb, note, 13, GREY, first=True, after=0)
    _footer(s)
    return s


def arch_slide():
    s = prs.slides.add_slide(BLANK)
    _header(s, "How information flows, end to end", "Deliverable 1 · Architecture")
    # data band
    _rect(s, 0.7, 1.62, 11.93, 0.5, NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
    db = _tb(s, 0.9, 1.66, 11.6, 0.42, MSO_ANCHOR.MIDDLE)
    _para(db, "DATA  ·  7 index sleeves + VIX + 7 FRED macro series  ·  weekly Fri close, shift(1), macro lagged 4w (no look-ahead)",
          12, WHITE, bold=True, first=True, after=0)
    # three model boxes
    boxes = [
        ("M1  ·  SIDE", ["Trains: rule-based, fixed weights", "Infers: score → weekly top-3", "Output: signal ∈ {0, +1}"]),
        ("M2  ·  QUALITY", ["Trains: logistic on M1≠0 rows", "Infers: P(trade succeeds)", "Output: p_success ∈ (0,1)"]),
        ("M3  ·  SIZE", ["Not trained — a fixed rule", "binary / linear / ECDF", "Output: M3_size ∈ [0,1]"]),
    ]
    x = 0.7; w = 3.7; gap = 0.42
    for i, (t, lines) in enumerate(boxes):
        left = x + i * (w + gap)
        _rect(s, left, 2.45, w, 1.9, CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = _tb(s, left + 0.2, 2.58, w - 0.4, 1.7)
        _para(tf, t, 15, NAVY, bold=True, first=True, after=6)
        for ln in lines:
            _para(tf, ln, 11.5, GREY, after=2)
        if i < 2:
            ar = _tb(s, left + w + 0.02, 2.9, gap, 0.6, MSO_ANCHOR.MIDDLE)
            _para(ar, "→", 24, ACCENT, bold=True, first=True, after=0, align=PP_ALIGN.CENTER)
    # portfolio + backtest
    _rect(s, 0.7, 4.6, 5.85, 0.95, CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
    p1 = _tb(s, 0.9, 4.68, 5.5, 0.8)
    _para(p1, "PORTFOLIO", 13, NAVY, bold=True, first=True, after=3)
    _para(p1, "raw_w = signal · M3_size · budget → 25%/name cap → 100% gross → 12% vol target", 11, GREY, after=0)
    _rect(s, 6.78, 4.6, 5.85, 0.95, CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
    p2 = _tb(s, 6.98, 4.68, 5.5, 0.8)
    _para(p2, "BACKTEST", 13, NAVY, bold=True, first=True, after=3)
    _para(p2, "weekly returns · turnover · 5 bps cost · vs EW 1/7 and 60/40", 11, GREY, after=0)
    _rect(s, 0.6, 5.75, 12.13, 0.72, LIGHT)
    tk = _tb(s, 0.85, 5.8, 11.7, 0.64, MSO_ANCHOR.MIDDLE)
    _para(tk, "In one line:  M1 sets recall (finds moves) · M2 sets precision (which pay) · M3 sets exposure (how much). "
              "M3 is the only stage that is a fixed function, not a model.", 12.5, ACCENT, bold=True, first=True, after=0)
    _footer(s)
    return s


# ================================================================ BUILD
title_slide(
    "Multi-Asset Meta-Labeling: architecture, evidence,\nand an honest read on M3",
    "Answering the four deliverables — with reproducible, walk-forward-controlled results.",
    "Ela Kumuk  ·  SSGA Field Project  ·  Research only, not investment advice",
)

bullets_slide(
    "Four deliverables, one authoritative run", "What this deck delivers",
    [
        ("Deliverable 1 — A complete architecture diagram: the inputs, outputs, training and inference of M1, M2, M3 and the portfolio layer.", 0),
        ("Deliverable 2 — A worked numerical example following one asset (SP500) through every stage to its final portfolio weight.", 0),
        ("Deliverable 3 — A structured experimental matrix spanning M1 weights, asset allocations, M3 sizing functions, regime conditioning and cost-awareness.", 0),
        ("Deliverable 4 — Exposure-controlled and walk-forward results that separate genuine model value from mechanical risk reduction.", 0),
        ("Plus the two conceptual asks: a precise account of meta-labeling, and one reproducible set of results where every number traces to a config, data window and code version.", 0),
    ],
    takeaway="This is treated as the beginning of the investigation, not a final verdict on meta-labeling or M3.",
)

arch_slide()

table_slide(
    "Each layer: what it trains on, what it infers", "Deliverable 1 · Architecture",
    ["Layer", "Trains on", "Infers (inference)", "Output"],
    [
        ["M1 · side", "Nothing — fixed economic weights", "Weekly cross-sectional score, take top-3", "signal ∈ {0, +1}"],
        ["M2 · quality", "Only M1≠0 weeks; label = did the trade pay", "P(this M1 trade beats the cost hurdle)", "p_success ∈ (0,1)"],
        ["M3 · size", "Not trained (deterministic rule)", "Map p_success → bet fraction", "M3_size ∈ [0,1]"],
        ["Portfolio", "—", "Caps, gross limit, volatility target", "final weights"],
    ],
    note="Key point for the reviewer: M1 and M3 involve no fitting. Only M2 is a trained classifier — and it is trained "
         "exclusively on the trades M1 proposed, which is what makes this meta-labeling rather than a second directional model.",
)

bullets_slide(
    "Meta-labeling, precisely — M2 judges M1's trade, it does not trade", "Conceptual clarity",
    [
        ("Which observations enter M2's training sample? Only asset-weeks where M1_signal ≠ 0 — the trades M1 actually proposed. M2 never sees flat weeks.", 0),
        ("How is the meta-label built? meta_label = 1 if sign(M1) · (forward 4-week return) > cost hurdle, else 0. It is a binary 'did this trade pay?', NOT a −1/+1 direction.", 0),
        ("Why a 4-week horizon and ±0.5% threshold? Four weeks matches the weekly-rebalance holding period; the ±0.5% band plus a 0.1% cost hurdle embeds a minimal edge above noise and frictions.", 0),
        ("Which features enter M2, and why? 52 features — momentum/trend/vol, cross-asset dispersion, macro regime, M1 context. Economically they answer: when does a momentum trade fail? (stress, crowding, adverse regime).", 0),
        ("How is the probability calibrated and validated? Sigmoid-calibrated; evaluated on the out-of-sample window only via AUC, calibration curve, and decile hit-rates.", 0),
    ],
    takeaway="M2 estimates the probability an M1 candidate succeeds under a defined outcome, horizon and cost — it is a quality filter, not a new signal.",
)

table_slide(
    "M2, read honestly", "Meta-labeling · validation",
    ["Metric", "Value", "Interpretation"],
    [
        ["Test AUC-ROC (full sample)", "0.589", "Weak but common for noisy financial labels"],
        ["Mean AUC across walk-forward folds", "≈ 0.49", "At/below random out-of-sample — the real test"],
        ["Recall @ threshold 0.55", "0.84", "Approves almost every trade → thin filter"],
        ["Mean P(winners) vs P(losers)", "0.577 vs 0.574", "Barely separates winners from losers"],
        ["Base rate (M1 trades that pay)", "60.8%", "The unconditional benchmark M2 must beat"],
    ],
    note="Honest read: as currently wired, M2's value is a SOFT input to M3 sizing, not a hard filter. Slide 12 shows this "
         "improves once macro is removed from M1 so M2 carries information M1 does not.",
    hi_row=2,
)

bullets_slide(
    "M3 — three rules map probability to a bet fraction", "M3 · sizing",
    [
        ("Binary:  f = 1[p ≥ T].  All-or-nothing at T = 0.55. Because recall ≈ 1, it approves nearly everything, so it behaves like M1-only.", 0),
        ("Linear:  f = max(0, 2p − 1).  Very defensive — high Sharpe but starves return. A drawdown tool, not a return engine.", 0),
        ("ECDF (production):  f = rank of p in the training distribution. Best risk-adjusted of the three in-sample.", 0),
        ("The critical caveat: the three rules differ mostly in HOW MUCH they deploy — not in which name is picked (M1 already picked). So a lower drawdown may be mechanical, not skill. Deliverable 4 tests exactly this.", 0),
    ],
    takeaway="M1 = opportunity selector · M2 = trade-quality probability · M3 = capital deployment · Portfolio = risk budget.",
)

table_slide(
    "One asset, every stage, to a final weight", "Deliverable 2 · Worked example (SP500, one week)",
    ["Stage", "Quantity", "Value", "Formula (from code)"],
    [
        ["Features", "mom / trend / macro / risk", "0.80 / 0.50 / 0.30 / 0.40", "z-scored, no look-ahead"],
        ["M1 score", "M1_score", "0.505", "0.45·.80 + 0.25·.50 + 0.20·.30 − 0.10·.40"],
        ["M1 signal", "rank 2 of 7 → top-3", "+1", "allocation_mode top_k, k=3"],
        ["M2", "p_success", "0.58", "calibrated logistic, 52 features"],
        ["M3 size", "binary / linear / ECDF", "1.00 / 0.16 / 0.62", "1[p≥.55] / max(0,2p−1) / ECDF(p)"],
        ["Raw weight", "raw_w (ECDF)", "0.0886", "1 × 0.62 × (1/7 budget)"],
        ["Vol target", "scale to 12% ann.", "× 1.36", "0.12 / 0.088 realized vol"],
        ["Final weight", "SP500 book weight", "≈ 12.0%", "0.0886 × 1.36"],
    ],
    note="Illustrative feature inputs; every transform is the exact production formula. For the SAME p_success, the three M3 "
         "rules give ~19% / ~3% / ~12% — they differ in deployment, not in which name was chosen.",
    hi_row=8,
    hi_col_map={(8, 2): AMBER, (2, 2): ACCENT, (4, 2): ACCENT},
)

table_slide(
    "A broad, staged search — not a fishing trip", "Deliverable 3 · Experimental matrix",
    ["Axis", "Levels to test", "Baseline"],
    [
        ["A · M1 factor weights", "merge mom+trend → technical; sweep technical/macro/risk", "45/25/20/10"],
        ["B · Selection rule", "top-K ∈ {2,3,4} · threshold · min_score", "top-3"],
        ["C · Asset allocation", "equal · score · inverse-vol · risk-parity · min-var · mean-var", "1/7 equal"],
        ["D · M3 sizing function", "binary(T-grid) · linear · ECDF · power · sigmoid · frac-Kelly · μ/σ", "ECDF"],
        ["E · Conditioning", "global · per-asset · regime (HMM state)", "global"],
        ["F · Cost-awareness", "none · turnover-penalized · no-trade band", "none"],
        ["G · Covariance", "diagonal · Ledoit-Wolf shrinkage", "diagonal"],
        ["H · Exposure control", "raw · gross-matched · vol-matched · constant-haircut", "→ D4"],
        ["I · Evaluation", "full-sample · 6-fold walk-forward · crisis periods · Deflated Sharpe", "walk-fwd"],
    ],
    note="Protocol: sweep one axis at a time, rank on WALK-FORWARD mean Sharpe, report every candidate with its Deflated "
         "Sharpe Ratio, then combine only walk-forward-positive winners. Axes A and (macro) are already run — next two slides.",
    hi_row=1,
)

table_slide(
    "Genuine value vs mechanical risk reduction", "Deliverable 4 · Exposure control",
    ["Fold", "M1-only Sharpe", "ECDF edge vs M1", "IR vs equal-weight"],
    [
        ["2015–16", "−0.13", "−0.38", "−0.53"],
        ["2017–18", "0.95", "−0.70", "−0.58"],
        ["2019–20", "0.57", "−0.14", "−0.98"],
        ["2021–22", "−0.06", "+0.92", "+0.55"],
        ["2023–24", "1.09", "−0.46", "−1.14"],
        ["2025–26", "1.83", "−0.38", "−1.08"],
        ["MEAN", "0.71", "−0.19", "−0.63"],
    ],
    note="ECDF beats M1-only in 1 of 6 folds; its lower drawdown coincides with ~52% gross vs EW's 100% — consistent with "
         "mechanical de-risking. The decisive tests: ① constant-haircut control, ② vol-matched comparison, ③ Spearman(M3_size, "
         "realized PnL). If ③ ≈ 0, M3 does not distinguish strong from weak trades.",
    hi_row=7,
    hi_col_map={(4, 2): POS, (4, 3): POS, (1, 2): NEG, (2, 2): NEG, (3, 2): NEG, (5, 2): NEG, (6, 2): NEG, (7, 2): NEG},
)

table_slide(
    "Result 1 — M1 factor weights were never tuned. We swept them.", "Broader experimentation · Axis A",
    ["M1 weights (mom/trend/macro/risk)", "Walk-fwd Sharpe", "Folds positive"],
    [
        ["risk_heavy   0.40/0.22/0.18/0.20", "0.811", "6 / 6"],
        ["baseline     0.45/0.25/0.20/0.10", "0.709", "4 / 6"],
        ["no_macro     0.56/0.31/0.00/0.13", "0.690", "4 / 6"],
        ["technical    0.35/0.35/0.20/0.10", "0.672", "4 / 6"],
        ["trend_only   0.00/0.70/0.20/0.10", "0.521", "5 / 6"],
        ["momentum_only 0.70/0.00/0.20/0.10", "0.500", "4 / 6"],
    ],
    note="Two findings. (1) Doubling the downside/risk-penalty weight lifts Sharpe 0.709 → 0.811 and is positive in ALL SIX "
         "folds — the most robust improvement found. (2) Momentum & trend are 0.78 correlated, yet either ALONE (~0.50) is far "
         "worse than both together (0.71): correlated is not redundant — the ensemble beats the parts.",
    hi_row=1,
    hi_col_map={(1, 1): POS, (1, 2): POS},
)

table_slide(
    "Result 2 — Is macro the lever? Two controlled tests.", "Broader experimentation · macro",
    ["Test", "Setup", "Finding"],
    [
        ["Macro in M1", "Hold sample fixed, vary which macro feeds M1", "Best set beats macro-OFF by only +0.013 Sharpe (t≈0.19) — noise"],
        ["Rates/credit in M1", "curve + credit signals into M1's tilt", "Actively hurts: 0.668 vs 0.701 with macro off"],
        ["Macro in M2 (M1 clean)", "Strip macro from M1, put it only in M2", "Two sets flip POSITIVE — drop-CPI edge +0.067, AUC 0.525"],
        ["Winner", "M2 macro = all series except CPI", "ECDF ≈ 0.745 > original M1-only 0.709 — meta-labeling finally adds"],
    ],
    note="Read: macro is NOT a lever for M1's selection (momentum/trend/VIX carry it). But once M1 is a clean technical model, "
         "the RIGHT macro set in M2 (drop backward-looking CPI) turns the meta-label positive — the textbook separation working. "
         "Caveat: edges are small and only 2/6 folds positive; Deflated-Sharpe and exposure-matching still required.",
    hi_row=3,
    hi_col_map={(3, 2): POS, (4, 2): POS, (1, 2): NEG, (2, 2): NEG},
)

bullets_slide(
    "One authoritative, reproducible set of results", "Reproducibility",
    [
        ("Every reported number traces to a specific config, data window, git SHA and timestamped run directory.", 0),
        ("The baseline walk-forward was reproduced bit-for-bit: same six folds, same −0.19 ECDF edge — the environment is trustworthy.", 0),
        ("A contradiction was found and flagged: an older report claimed an ECDF edge of +0.177 across 4-of-6 folds; that was an ETF-era artifact. The true index-era value is −0.19 across 1-of-6.", 0),
        ("All experiments here use one canonical config with only the tested axis changed — so any difference is attributable to that axis, not a shifted sample.", 0),
    ],
    takeaway="No number without a run behind it. Reconciling every report to the reproduced index-era run is the deliverable.",
)

bullets_slide(
    "Honest verdict and next phase", "Where this stands",
    [
        ("What holds up: M1-only is the robust core (≈ EW return, ~half the drawdown, higher Sharpe); and raising the risk-penalty weight (Axis A) is a real, fold-consistent improvement to 0.811 Sharpe.", 0),
        ("What is still unproven: M2/M3 meta-labeling does not beat M1-only out-of-sample in most folds, and ECDF's drawdown edge may be mechanical until the exposure-matched tests confirm otherwise.", 0),
        ("Most promising thread: a clean meta-labeling separation (technical M1 + macro-in-M2 excluding CPI) is the first configuration where the meta-label adds value — worth deepening.", 0),
        ("Next: (1) exposure-matched M3 comparison, (2) regime-conditioned M3, (3) deepen the technical-M1 + macro-M2 design, (4) broaden the asset-allocation axis.", 0),
    ],
    takeaway="A rigorous baseline that names its own weaknesses — and defines the exact tests that will resolve them.",
    takeaway_color=NAVY,
)

out = Path(__file__).resolve().parents[1] / "presentation" / "SSGA_Deliverables_Deck.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print("saved:", out, "|", len(prs.slides._sldIdLst), "slides")
