"""Generate Summer 2026 presentation deck (Jun 28 – Aug 8).

    pip install -e ".[slides]"
    python scripts/make_summer_deck.py
        -> reports/Summer_2026_Meta_Labeling_Deck.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
GREY = RGBColor(0x55, 0x5B, 0x66)
ACCENT = RGBColor(0x1F, 0x6F, 0xB2)
LIGHT = RGBColor(0xEC, 0xF1, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DATE = "August 8, 2026"
FOOT = "SSGA Summer 2026 · Meta-Labeling Pipeline · " + DATE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_n = [0]


def _tb(slide, left, top, width, height, anchor=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def _rect(slide, left, top, width, height, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _header(slide, title):
    _rect(slide, 0.0, 0.0, 0.22, 7.5, NAVY)
    t = _tb(slide, 0.6, 0.35, 12.4, 0.9)
    p = t.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY
    _rect(slide, 0.62, 1.18, 4.2, 0.045, ACCENT)


def _footer(slide):
    _n[0] += 1
    f = _tb(slide, 0.6, 7.02, 12.2, 0.35)
    p = f.paragraphs[0]
    p.text = FOOT
    p.font.size = Pt(9)
    p.font.color.rgb = GREY
    pg = _tb(slide, 12.4, 7.02, 0.7, 0.35)
    p = pg.paragraphs[0]
    p.text = str(_n[0])
    p.font.size = Pt(9)
    p.font.color.rgb = GREY


def title_slide(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0.0, 0.0, 13.333, 0.35, NAVY)
    _rect(s, 0.0, 7.15, 13.333, 0.35, ACCENT)
    tf = _tb(s, 0.9, 2.2, 11.5, 2.8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = GREY
    p2.space_before = Pt(10)
    p3 = tf.add_paragraph()
    p3.text = DATE
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT
    p3.space_before = Pt(8)
    return s


def content(title, bullets, takeaway=None):
    s = prs.slides.add_slide(BLANK)
    _header(s, title)
    body = _tb(s, 0.7, 1.45, 12.0, 4.7)
    for i, (txt, lvl) in enumerate(bullets):
        par = body.paragraphs[0] if i == 0 else body.add_paragraph()
        par.text = ("■  " if lvl == 0 else "–  ") + txt
        par.level = lvl
        par.font.size = Pt(17 if lvl == 0 else 14)
        par.font.bold = lvl == 0
        par.font.color.rgb = NAVY if lvl == 0 else GREY
        par.space_after = Pt(5)
    if takeaway:
        _rect(s, 0.6, 6.25, 12.13, 0.62, LIGHT)
        tk = _tb(s, 0.8, 6.3, 11.8, 0.55, MSO_ANCHOR.MIDDLE)
        p = tk.paragraphs[0]
        p.text = "Takeaway:  " + takeaway
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT
    _footer(s)
    return s


def table_slide(title, headers, rows, note=None, highlight_row=None, note_label="Note"):
    s = prs.slides.add_slide(BLANK)
    _header(s, title)
    nr, nc = len(rows) + 1, len(headers)
    rh = min(0.38, 4.5 / max(nr, 1))
    gt = s.shapes.add_table(nr, nc, Inches(0.7), Inches(1.5), Inches(11.9), Inches(rh * nr)).table
    for c, h in enumerate(headers):
        cell = gt.cell(0, c)
        cell.text = h
        pr = cell.text_frame.paragraphs[0]
        pr.font.size = Pt(12)
        pr.font.bold = True
        pr.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows, start=1):
        hot = highlight_row is not None and r - 1 == highlight_row
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.text = str(val)
            pr = cell.text_frame.paragraphs[0]
            pr.font.size = Pt(12)
            pr.font.bold = hot
            pr.font.color.rgb = ACCENT if hot else NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if hot else WHITE
    if note:
        top = 1.5 + rh * nr + 0.2
        _rect(s, 0.6, top, 12.13, 0.95, LIGHT)
        tk = _tb(s, 0.8, top + 0.06, 11.8, 0.85, MSO_ANCHOR.MIDDLE)
        p = tk.paragraphs[0]
        p.text = note_label + ":  " + note
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = ACCENT
    _footer(s)
    return s


# --- Slides (Jun 28 – Aug 8, 2026) ---

title_slide(
    "Multi-Asset Meta-Labeling Pipeline",
    "Summer 2026 update · Jun 28 – Aug 8 · research & educational use only",
)

content(
    "Timeline — What We Did",
    [
        ("Jun 28 – early Jul: simplify M1 (static factors), LR-only M2, macro→M2, index signals.", 0),
        ("Jul: Joubert M3 formalization; factor/regime/M3 diagnostics; 6+ companion reports.", 0),
        ("Jul–Aug research: walk-forward, M1 weight tuning, M3 threshold sweep, IR attribution, M2 enrichment.", 0),
        ("Aug 8: index-first sleeve IDs (SP500, …); full validation; docs + deck alignment.", 0),
    ],
    takeaway="Explainability and stability gates matter as much as holdout Sharpe.",
)

content(
    "Pipeline Architecture (Joubert)",
    [
        ("M1 — which side? Top-3 long index sleeves per week (rule-based).", 0),
        ("M2 — P(trade succeeds)? Logistic regression, 52 features, calibrated.", 0),
        ("M3 — how much to bet? ECDF sizing maps p_success → M3_size ∈ [0,1].", 0),
        ("Portfolio — per-sleeve caps, gross exposure, 12% vol target, 5 bps costs.", 0),
    ],
    takeaway="M2 ranks quality; M3 shapes risk — not a second classifier at T=0.55.",
)

table_slide(
    "Index Universe (Aug 2026)",
    ["Sleeve ID", "Asset class", "Free data source", "Notes"],
    [
        ["SP500", "U.S. equity", "^GSPC", "true index"],
        ["MSCI_EAFE", "Dev. ex-US", "EFA proxy", "ETF proxy"],
        ["MSCI_EM", "EM equity", "EEM proxy", "ETF proxy"],
        ["UST_7_10", "Govt bonds", "IEF proxy", "ETF proxy"],
        ["US_HIGH_YIELD", "Credit", "HYG proxy", "ETF proxy"],
        ["GOLD_SPOT", "Gold", "GC=F", "futures proxy"],
        ["US_REIT", "Real estate", "FRED REIT", "true index · binds panel ~2011"],
    ],
    note="Signals on index/proxy series; sleeve IDs are research labels, not trade tickets.",
)

table_slide(
    "OOS Results — Long-Only (Test 2021+)",
    ["Strategy", "Ann. Return", "Sharpe", "Max DD"],
    [
        ["Equal Weight 1/7", "7.81%", "0.74", "-22.3%"],
        ["M1 Only", "8.92%", "0.88", "-20.2%"],
        ["M1+M2+M3 Binary", "8.67%", "0.95", "-14.4%"],
        ["M1+M2+M3 ECDF", "4.58%", "0.91", "-7.5%"],
    ],
    note="From latest index-sleeve pipeline run — see reports/final_report.md.",
    highlight_row=2,
)

content(
    "M1 — Factor Model",
    [
        ("Weights: momentum 45%, trend 25%, macro 20%, risk penalty 10%.", 0),
        ("Rank all 7 sleeves weekly; select top K=3 longs.", 0),
        ("IC-proportional weight tuning looked better on holdout but failed walk-forward.", 0),
        ("M1-only economics stable vs prior branch on comparable windows.", 0),
    ],
    takeaway="M1 selects opportunities; keep baseline 45/25/20/10.",
)

content(
    "M2 — Meta-Label",
    [
        ("52 enriched features (M1 components, interactions, asset-class dummies).", 0),
        ("Test AUC ~0.59 vs ~0.57 baseline (+0.02).", 0),
        ("Trees / per-asset heads overfit — global LR retained.", 0),
        ("m1_components_rich variant rejected on walk-forward (2/6 fold wins).", 0),
    ],
    takeaway="Modest ranking lift; main value is feeding M3 ECDF.",
)

content(
    "M3 — Bet Sizing & Allocation States",
    [
        ("no_signal — M1 flat; m3_zero — candidate vetoed; m3_active — positive size.", 0),
        ("Binary T=0.55 approves ~100% of candidates (degenerate vs M1-only).", 0),
        ("ECDF production mode: mean M3_size ~0.53 on candidates; rarely exact zero.", 0),
        ("M3_size is magnitude; direction from M1_signal (±1).", 0),
    ],
    takeaway="ECDF is the production risk-shaping layer.",
)

content(
    "m3_zero — Long-Only vs Long/Short",
    [
        ("Long-only: m3_zero = M1 wanted a long, M3 allocated zero.", 0),
        ("Long/short: same rule for long OR short candidates.", 0),
        ("ECDF long-only: ~0% m3_zero; binary/linear long/short: high rejection.", 0),
        ("Distinct from no_signal (M1 never nominated the sleeve).", 0),
    ],
    takeaway="M3 veto is pre-portfolio-cap; not the same as M1=0.",
)

content(
    "Walk-Forward & Research Verdicts",
    [
        ("6 expanding-window folds; ECDF edge positive in 4/6 vs M1-only.", 0),
        ("Rejected: IC M1 weights, M3 threshold promotion, IR overlays, M2 enrichment.", 0),
        ("IR vs Sharpe trade-off: ECDF under-invests vs equal-weight (~52% gross).", 0),
        ("74 pytest tests + integration smoke on index sleeves.", 0),
    ],
    takeaway="Holdout wins require walk-forward confirmation before config change.",
)

content(
    "Open Directions",
    [
        ("Regime-conditioned M3 sizing (weak 2015–16 fold).", 0),
        ("Promote M3 thresholds only with walk-forward gates.", 0),
        ("Long/short sleeve still weak vs long-only.", 0),
        ("Bloomberg true-index history for pre-2011 research.", 0),
    ],
    takeaway="Next high-value work: regime-aware M3, not more M2 features without stability proof.",
)

content(
    "Disclaimer",
    [
        ("Research backtest on public index/proxy data.", 0),
        ("Not investment advice; no live execution or capacity modeling.", 0),
        ("See TERMINOLOGY.md and DATA_SOURCES_AND_ETL.md for definitions.", 0),
    ],
)

root = Path(__file__).resolve().parent.parent
out = root / "reports" / "Summer_2026_Meta_Labeling_Deck.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"saved {_n[0] + 1} slides -> {out}")
